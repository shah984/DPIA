"""
document_reader.py — Convert uploaded files to plain text.

Supports: .pdf (via pdfplumber), .docx (via python-docx), .pptx (via python-pptx), .txt
Returns empty string on unreadable files (caller decides how to surface the warning).
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def read_file(path: str | Path) -> tuple[str, str | None]:
    """
    Read a file and return its plain-text content.

    Returns:
        (text, error_message) — error_message is None on success, a warning string on failure.
    """
    path = Path(path)
    suffix = path.suffix.lower()

    try:
        if suffix == ".pdf":
            return _read_pdf(path), None
        elif suffix == ".docx":
            return _read_docx(path), None
        elif suffix == ".pptx":
            return _read_pptx(path), None
        elif suffix == ".txt":
            return _read_txt(path), None
        else:
            return "", f"Unsupported file type '{suffix}' — skipped."
    except Exception as exc:
        logger.warning("Failed to read %s: %s", path.name, exc)
        return "", f"Could not read '{path.name}': {exc}"


def _read_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text.strip())
    return "\n\n".join(parts)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append("  |  ".join(cells))

    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def _read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def read_all(paths: list[str | Path]) -> tuple[str, list[str]]:
    """
    Read multiple files and concatenate their text.

    Returns:
        (combined_text, warnings) — warnings is a list of error strings for failed files.
    """
    all_parts: list[str] = []
    warnings: list[str] = []

    for p in paths:
        text, err = read_file(p)
        if err:
            warnings.append(err)
        elif text.strip():
            name = Path(p).name
            all_parts.append(f"=== {name} ===\n{text}")

    return "\n\n".join(all_parts), warnings
